#!/usr/bin/env python3
"""Build the RCN / CMTAT / FireFly gateway deck (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# palette
INK   = RGBColor(0x0F, 0x1A, 0x2B)  # near-black navy
NAVY  = RGBColor(0x14, 0x2A, 0x4A)
ACCENT= RGBColor(0x2D, 0x6C, 0xDF)  # blue
TEAL  = RGBColor(0x11, 0x9D, 0xA4)
GOLD  = RGBColor(0xC8, 0x92, 0x2B)
MUTE  = RGBColor(0x5A, 0x67, 0x78)
LINE  = RGBColor(0xD9, 0xDF, 0xE7)
BG    = RGBColor(0xF6, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s

def box(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    b = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    b.shadow.inherit = False
    if fill is None: b.fill.background()
    else: b.fill.solid(); b.fill.fore_color.rgb = fill
    if line is None: b.line.fill.background()
    else: b.line.color.rgb = line; b.line.width = Pt(lw)
    return b

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=1.0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs[0], tuple): runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = sp
        for (t, sz, col, bold) in para:
            r = p.add_run(); r.text = t; f = r.font
            f.size = Pt(sz); f.color.rgb = col; f.bold = bold; f.name = "Calibri"
    return tb

def bullet(s, x, y, w, items, sz=15, gap=None, col=INK, mk=ACCENT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(4))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=0; tf.margin_top=0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap if gap is not None else 8); p.line_spacing = 1.05
        r = p.add_run(); r.text = "▸ "; r.font.size = Pt(sz); r.font.color.rgb = mk; r.font.bold = True
        if isinstance(it, tuple):
            head, rest = it
            r2 = p.add_run(); r2.text = head; r2.font.size=Pt(sz); r2.font.color.rgb=col; r2.font.bold=True
            r3 = p.add_run(); r3.text = rest; r3.font.size=Pt(sz); r3.font.color.rgb=MUTE; r3.font.bold=False
        else:
            r2 = p.add_run(); r2.text = it; r2.font.size=Pt(sz); r2.font.color.rgb=col
        for rr in p.runs: rr.font.name = "Calibri"
    return tb

def kicker(s, text, x=0.9, y=0.62):
    box(s, x, y+0.02, 0.055, 0.34, fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
    txt(s, x+0.18, y, 9, 0.4, [[(text.upper(), 13, GOLD, True)]])

def header(s, title, kick=None):
    if kick: kicker(s, kick)
    txt(s, 0.9, 0.95, 11.6, 0.9, [[(title, 30, INK, True)]])
    box(s, 0.9, 1.72, 11.53, 0.02, fill=LINE, shape=MSO_SHAPE.RECTANGLE)

def pagenum(s, n):
    txt(s, 12.3, 7.02, 0.9, 0.3, [[(f"{n:02d}", 11, MUTE, False)]], align=PP_ALIGN.RIGHT)

# ---------- 1 TITLE ----------
s = slide(INK)
box(s, 0, 0, SW, SH, fill=INK, shape=MSO_SHAPE.RECTANGLE)
box(s, 0, 6.9, SW, 0.6, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
box(s, 0.9, 1.5, 0.09, 1.2, fill=GOLD, shape=MSO_SHAPE.RECTANGLE)
txt(s, 1.2, 1.35, 11, 0.5, [[("REFERENCE ARCHITECTURE", 15, GOLD, True)]])
txt(s, 1.2, 1.95, 11.2, 2.4, [
    [("Instant-Payment Front-to-Back", 42, WHITE, True)],
    [("Settlement Gateway for Tokenized RCNs", 42, WHITE, True)],
], sp=1.02)
txt(s, 1.2, 4.35, 11.2, 1.3, [
    [("Cash flows ", 19, RGBColor(0xBF,0xD2,0xF0), False),("TradFi → TradFi", 19, GOLD, True),(".  The permissioned DLT is the shared", 19, RGBColor(0xBF,0xD2,0xF0), False)],
    [("messaging / coordination layer — FireFly minting CMTAT tokens for RCNs.", 19, RGBColor(0xBF,0xD2,0xF0), False)],
], sp=1.1)
txt(s, 1.2, 6.98, 11, 0.4, [[("DLT engineering · FireFly integration · cross-border settlement · multi-jurisdiction compliance", 12, RGBColor(0x9A,0xAD,0xC8), False)]])

# ---------- 2 THE PRODUCT ----------
s = slide(); header(s, "The instrument — Reverse Convertible Note", "What we tokenize"); pagenum(s,2)
box(s, 0.9, 2.05, 5.5, 1.15, fill=WHITE, line=LINE)
txt(s, 1.15, 2.2, 5.1, 1.0, [
    [("RCN = ", 17, INK, True),("zero-coupon note", 17, ACCENT, True),("  +  ", 17, INK, True),("short put", 17, TEAL, True)],
    [("Above-market coupon funded by the option premium the investor implicitly sells.", 12.5, MUTE, False)],
], anchor=MSO_ANCHOR.MIDDLE, sp=1.1)
bullet(s, 0.9, 3.5, 5.5, [
    ("Barrier ≥ level → ", "cash redemption at par + coupon (token burned)"),
    ("Barrier < level → ", "physical delivery of underlying, at a loss (token burned)"),
    ("Variants: ", "barrier / knock-in, worst-of basket, autocallable"),
], sz=14.5)
box(s, 6.75, 2.05, 5.68, 4.6, fill=RGBColor(0xEE,0xF3,0xFB), line=LINE)
txt(s, 7.05, 2.25, 5.1, 0.5, [[("Why classification matters", 16, NAVY, True)]])
bullet(s, 7.05, 2.9, 5.1, [
    ("Structured DEBT security", " — not e-money, not a payment token."),
    ("Drives the overlay:", " prospectus / PRIIPs KID, MiFID II product governance & appropriateness."),
    ("Transfer-agent duties", " and investor-eligibility gating apply."),
    ("Token is the register", " where DLT law allows (ledger-based security)."),
], sz=14)

# ---------- 3 WHY CMTAT + FIREFLY ----------
s = slide(); header(s, "Why CMTAT + Hyperledger FireFly", "Building blocks"); pagenum(s,3)
box(s, 0.9, 2.0, 5.55, 4.7, fill=WHITE, line=LINE)
txt(s, 1.15, 2.18, 5.1, 0.5, [[("CMTAT", 19, ACCENT, True),("  security-token framework", 14, MUTE, False)]])
bullet(s, 1.15, 2.85, 5.05, [
    ("Mint / Burn", " — issue on cash finality; redeem at maturity"),
    ("Pause / Enforcement", " — halt issuance; freeze a holder"),
    ("RuleEngine (ERC-1404)", " — eligibility & jurisdiction gating"),
    ("Snapshot", " — holders-of-record at each coupon date"),
    ("Debt module", " — coupon schedule, maturity, credit events"),
], sz=13.5)
box(s, 6.78, 2.0, 5.65, 4.7, fill=WHITE, line=LINE)
txt(s, 7.03, 2.18, 5.1, 0.5, [[("FireFly", 19, TEAL, True),("  orchestration supernode", 14, MUTE, False)]])
bullet(s, 7.03, 2.85, 5.15, [
    ("Token connector", " — business-level mint/transfer/burn API"),
    ("Blockchain connector", " — nonce, gas, retry, receipts"),
    ("Event streams", " — durable, at-least-once, checkpointed"),
    ("Data exchange", " — PII / KID off-chain, hash on-chain"),
    ("Transaction manager", " — idempotent, exactly-once mint"),
], sz=13.5, mk=TEAL)
txt(s, 0.9, 6.85, 11.5, 0.4, [[("Net: FireFly turns “cash settled with finality” into “token minted, once, to the right wallet, with the right restrictions.”", 13, NAVY, True)]])

# ---------- 4 FRONT-TO-BACK ARCH ----------
s = slide(); header(s, "Front-to-back architecture", "How it fits together"); pagenum(s,4)
def band(x, w, title, col):
    box(s, x, 2.05, w, 4.5, fill=WHITE, line=LINE)
    box(s, x, 2.05, w, 0.5, fill=col, shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s, x, 2.11, w, 0.4, [[(title, 13.5, WHITE, True)]], align=PP_ALIGN.CENTER)
band(0.9, 3.7, "① CASH — TradFi → TradFi", NAVY)
band(4.82, 3.7, "② GATEWAY — correlates", ACCENT)
band(8.74, 3.7, "③ DLT — MESSAGING (no cash)", TEAL)
def node(x, y, w, t, sub, col):
    box(s, x, y, w, 0.72, fill=BG, line=col, lw=1.25)
    txt(s, x+0.1, y+0.08, w-0.2, 0.6, [
        [(t, 12.5, INK, True)],[(sub, 10, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE, sp=1.0)
node(1.05, 2.8, 3.4, "Client / Investor", "fiat in ↓   ·   fiat out ↑", NAVY)
node(1.05, 3.62, 3.4, "Instant Rail IN", "subscribe — SEPA Inst · FedNow · RTP", NAVY)
node(1.05, 4.44, 3.4, "Settlement / Paying Bank", "escrow · client-money", NAVY)
node(1.05, 5.26, 3.4, "Instant Rail OUT", "coupon · redemption par → client", NAVY)
node(4.97, 2.8, 3.4, "Payment Gateway", "ISO 20022 ingestion", ACCENT)
node(4.97, 3.62, 3.4, "Recon & Idempotency", "source of truth", ACCENT)
node(4.97, 4.44, 3.4, "Compliance Engine", "KYC/AML · Travel Rule · eligibility", ACCENT)
node(4.97, 5.26, 3.4, "FireFly Orchestrator", "mint / transfer / burn", ACCENT)
node(8.89, 2.8, 3.4, "CMTAT RCN Token", "shared golden record + message", TEAL)
node(8.89, 3.62, 3.4, "Authoritative Register", "ledger-based security", TEAL)
node(8.89, 4.44, 3.4, "Investor Wallet", "whitelisted · custodial / MPC", TEAL)
node(8.89, 5.26, 3.4, "Coordination Events", "mint · snapshot · burn → payout", TEAL)

# ---------- 5 ATOMICITY ----------
s = slide(); header(s, "The coordination problem — the actual hard part", "No on-chain value leg to make atomic"); pagenum(s,5)
txt(s, 0.9, 1.85, 11.5, 0.5, [[("Cash stays in TradFi; the DLT only messages. So the problem isn’t atomic value transfer — it’s correlating an off-chain cash event to an on-chain message, exactly once, with a safe reversal.", 13.5, MUTE, False)]])
rows = [
    ("Naïve sequential", "Cash finality → then message (mint)", "Window of “cash taken, note not recorded.” Must reconcile + reverse.", MUTE),
    ("Escrow + msg-on-finality  ✅", "Escrow held; mint only on confirmed finality; auto-refund on timeout", "Strong client protection; needs safeguarded account. Recommended.", ACCENT),
    ("Escrow release gated on mint", "Escrow → issuer only AFTER on-chain mint event observed", "Kills both races: minted-but-cash-lost and cash-released-but-not-minted.", TEAL),
    ("Optional on-chain DvP", "Tokenized deposit on same ledger, where law permits", "Cleanest atomicity — but pulls money on-chain, against the TradFi→TradFi goal.", GOLD),
]
y = 2.55
box(s, 0.9, y, 11.53, 0.42, fill=INK, shape=MSO_SHAPE.RECTANGLE)
for lbl,x,w in [("MODEL",1.05,2.7),("MECHANISM",3.9,4.2),("TRADE-OFF",8.2,4.1)]:
    txt(s, x, y+0.06, w, 0.3, [[(lbl,11.5,WHITE,True)]])
y += 0.42
for name, mech, trade, col in rows:
    box(s, 0.9, y, 11.53, 0.86, fill=WHITE, line=LINE)
    box(s, 0.9, y, 0.09, 0.86, fill=col, shape=MSO_SHAPE.RECTANGLE)
    txt(s, 1.05, y+0.1, 2.75, 0.7, [[(name,13,INK,True)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 3.9, y+0.1, 4.2, 0.7, [[(mech,11.5,MUTE,False)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 8.2, y+0.1, 4.1, 0.7, [[(trade,11.5,MUTE,False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.86
txt(s, 0.9, y+0.08, 11.5, 0.4, [[("Recommended: ", 13, GOLD, True),("escrow + message-on-finality, escrow release gated on the mint event, idempotent reconciliation as source of truth. No value on-chain = tiny failure surface.", 13, INK, False)]])

# ---------- 6 SEQUENCE ----------
s = slide(); header(s, "Payment → mint sequence (escrow, idempotent)", "Exactly-once"); pagenum(s,6)
lanes = ["Client","Settlement\nBank","Payment\nGateway","Recon /\nIdempotency","FireFly","CMTAT\nToken"]
n = len(lanes); x0 = 1.1; span = 11.1; step = span/(n-1)
top = 2.35; bot = 6.5
for i,l in enumerate(lanes):
    cx = x0 + i*step
    box(s, cx-0.85, top-0.55, 1.7, 0.5, fill=NAVY)
    txt(s, cx-0.85, top-0.52, 1.7, 0.45, [[(l,10.5,WHITE,True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, sp=0.9)
    ln = box(s, cx-0.006, top, 0.012, bot-top, fill=LINE, shape=MSO_SHAPE.RECTANGLE)
def arrow(i,j,y,label,col=INK,dash=False):
    cx1=x0+i*step; cx2=x0+j*step
    c = s.shapes.add_connector(2, Inches(min(cx1,cx2)), Inches(y), Inches(max(cx1,cx2)), Inches(y))
    c.line.color.rgb=col; c.line.width=Pt(1.4)
    le = c.line._get_or_add_ln()
    end = "tailEnd" if cx1>cx2 else "headEnd"
    e = le.makeelement(qn('a:'+end), {'type':'triangle','w':'med','len':'med'}); le.append(e)
    if dash:
        d=le.makeelement(qn('a:prstDash'),{'val':'dash'}); le.insert(0,d)
    midx=(cx1+cx2)/2
    txt(s, midx-1.6, y-0.32, 3.2, 0.3, [[(label,10,col,False)]], align=PP_ALIGN.CENTER)
arrow(0,1,2.7,"instant credit transfer (pacs.008)")
arrow(1,2,3.15,"pacs.002 ACSC — settlement finality",TEAL)
arrow(2,3,3.6,"record(subscriptionId, amount)")
txt(s,x0+3*step-0.85,3.85,1.7,0.3,[[("idempotency check",9.5,MUTE,False)]],align=PP_ALIGN.CENTER)
arrow(3,4,4.35,"mintRequest [operationId = subscriptionId]",ACCENT)
arrow(4,5,4.8,"mint(wallet, units)",ACCENT)
arrow(5,4,5.25,"Mint event + receipt",TEAL,dash=True)
arrow(4,3,5.7,"minted(subscriptionId, txHash)",TEAL,dash=True)
arrow(3,1,6.15,"release escrow to issuer",GOLD)
box(s, 0.9, 6.72, 11.53, 0.5, fill=RGBColor(0xEE,0xF3,0xFB), line=LINE)
txt(s, 1.1, 6.8, 11.1, 0.35, [[("Idempotency key = ISO 20022 end-to-end id = FireFly operationId → re-delivered payment events never double-mint. On reject/timeout → auto-refund escrow.", 11.5, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE)

# ---------- 7 LIFECYCLE ----------
s = slide(); header(s, "CMTAT ↔ RCN lifecycle mapping", "Token state machine"); pagenum(s,7)
def st(x,y,w,t,col,fill=WHITE):
    box(s,x,y,w,0.68,fill=fill,line=col,lw=1.5)
    txt(s,x,y+0.06,w,0.56,[[(t,13,INK,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def conn(x1,y1,x2,y2,label,col=MUTE,dash=False):
    c=s.shapes.add_connector(2,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=col;c.line.width=Pt(1.6)
    le=c.line._get_or_add_ln()
    e=le.makeelement(qn('a:headEnd'),{'type':'triangle','w':'med','len':'med'});le.append(e)
    if dash:
        d=le.makeelement(qn('a:prstDash'),{'val':'dash'});le.insert(0,d)
    if label:
        txt(s,(x1+x2)/2-1.5,(y1+y2)/2-0.34,3.0,0.3,[[(label,9.5,col,False)]],align=PP_ALIGN.CENTER)
# linear spine
st(0.95,3.6,1.5,"Subscribed",NAVY)
st(3.15,3.6,1.7,"Active (in issue)",ACCENT,fill=RGBColor(0xEE,0xF3,0xFB))
conn(2.45,3.94,3.15,3.94,"mint()",ACCENT)
# choice diamond
dx,dy=6.55,3.55
dia=box(s,dx,dy,0.9,0.9,fill=GOLD,shape=MSO_SHAPE.DIAMOND)
txt(s,dx-0.4,dy+0.95,1.7,0.3,[[("barrier?",10,GOLD,True)]],align=PP_ALIGN.CENTER)
conn(4.85,3.94,6.55,4.0,"maturity / autocall",INK)
st(8.0,2.75,2.0,"Redeemed",TEAL)
st(8.0,4.6,2.0,"Settled (physical)",TEAL)
conn(7.45,3.85,8.0,3.15,"OK → burn + par",TEAL)
conn(7.45,4.15,8.0,4.85,"breached → deliver",GOLD)
st(11.0,2.75,1.7,"End",MUTE); st(11.0,4.6,1.7,"End",MUTE)
conn(10.0,3.09,11.0,3.09,"",MUTE); conn(10.0,4.94,11.0,4.94,"",MUTE)
# coupon loop
st(3.15,5.7,1.7,"CouponRecord",TEAL,fill=RGBColor(0xEB,0xF6,0xF6))
conn(4.0,4.28,4.0,5.7,"snapshot() @ record date",TEAL)
conn(3.15,5.9,2.6,5.9,"",TEAL)
conn(2.6,5.9,2.6,4.05,"coupon paid on fiat rail",TEAL,dash=True)
# frozen
st(3.15,1.85,1.7,"Frozen",RGBColor(0xC0,0x3A,0x3A),fill=RGBColor(0xFB,0xEE,0xEE))
conn(4.0,3.6,4.0,2.53,"freeze(holder)",RGBColor(0xC0,0x3A,0x3A))
box(s, 0.9, 6.75, 11.53, 0.5, fill=BG, line=LINE)
txt(s,1.05,6.82,11.2,0.36,[[("Three ways out of Active: ",11.5,INK,True),("recurring coupon cycle · involuntary Frozen hold (sanctions/court) · terminal Exit where the barrier routes to cash vs physical.",11.5,MUTE,False)]],anchor=MSO_ANCHOR.MIDDLE)

# ---------- 8 JURISDICTION MATRIX ----------
s = slide(); header(s, "Multi-jurisdiction compliance matrix", "Where cross-border bites"); pagenum(s,8)
cols = ["","🇨🇭 Switzerland","🇪🇺 EU","🇬🇧 UK","🇸🇬 Singapore","🇺🇸 US"]
data = [
 ("Instant rail","SIC / SEPA","SEPA Inst · TIPS","Faster Payments","FAST · PayNow","FedNow · RTP"),
 ("Legal basis","DLT Act (CO 973d)","MiCA + MiFID II","FSMA · DSS","SFA · PS Act","Securities Act"),
 ("Disclosure","FIDLEG KID","PRIIPs + Prospectus","UK PRIIPs · Duty","Product sheet","Reg S / 144A"),
 ("Investor gate","Qualified/retail","MiFID categories","FCA categorisation","Accredited/instl","Accredited/QIB"),
 ("AML / Travel","AMLA","AMLD · TFR","MLR 2017","MAS Notice","BSA / FinCEN"),
]
x0=0.9; wlabel=1.9; wc=(11.53-wlabel)/5; y=2.05; rh=0.82
# header row
box(s,x0,y,wlabel,rh,fill=INK,shape=MSO_SHAPE.RECTANGLE)
for j in range(5):
    box(s,x0+wlabel+j*wc,y,wc,rh,fill=NAVY,shape=MSO_SHAPE.RECTANGLE)
    txt(s,x0+wlabel+j*wc,y+0.14,wc,0.55,[[(cols[j+1],12.5,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sp=0.9)
y+=rh
for r,(lbl,*vals) in enumerate(data):
    fill = WHITE if r%2==0 else BG
    box(s,x0,y,wlabel,rh,fill=RGBColor(0x23,0x35,0x52),shape=MSO_SHAPE.RECTANGLE)
    txt(s,x0+0.1,y+0.1,wlabel-0.2,rh-0.2,[[(lbl,11.5,WHITE,True)]],anchor=MSO_ANCHOR.MIDDLE)
    for j,v in enumerate(vals):
        box(s,x0+wlabel+j*wc,y,wc,rh,fill=fill,line=LINE)
        txt(s,x0+wlabel+j*wc+0.06,y+0.08,wc-0.12,rh-0.16,[[(v,11,INK,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,sp=0.95)
    y+=rh
txt(s,0.9,y+0.12,11.5,0.5,[[("Enforced in the RuleEngine, not the UI: ",12.5,GOLD,True),("wallet jurisdiction is checked on every transfer · Travel-Rule payloads off-chain (IVMS-101) · PII hashes only.",12.5,INK,False)]])

# ---------- 9 CODE ----------
s = slide(INK);
box(s,0,0,SW,SH,fill=INK,shape=MSO_SHAPE.RECTANGLE)
kicker(s,"Implementation",x=0.9,y=0.55)
txt(s,1.08,0.9,11.6,0.8,[[("Two thin layers — token + orchestrator",28,WHITE,True)]])
box(s,0.9,1.66,11.53,0.02,fill=NAVY,shape=MSO_SHAPE.RECTANGLE)
def code(x,w,title,lines,col):
    box(s,x,2.0,w,4.7,fill=RGBColor(0x0A,0x12,0x1F),line=NAVY,lw=1.0)
    box(s,x,2.0,w,0.44,fill=col,shape=MSO_SHAPE.ROUND_2_SAME_RECTANGLE)
    txt(s,x+0.2,2.06,w-0.4,0.34,[[(title,12.5,WHITE,True)]])
    tb=s.shapes.add_textbox(Inches(x+0.18),Inches(2.6),Inches(w-0.34),Inches(4.0))
    tf=tb.text_frame; tf.word_wrap=True
    for i,(t,c) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing=1.12
        r=p.add_run(); r.text=t; r.font.size=Pt(10.5); r.font.name="Consolas"; r.font.color.rgb=c
CM=RGBColor(0x5F,0x8A,0xC0); KW=RGBColor(0xC8,0x92,0x2B); TXTC=RGBColor(0xD7,0xE1,0xF0); STR=RGBColor(0x6B,0xC5,0x8E)
code(0.9,5.6,"RCNToken.sol  (extends CMTAT)",[
 ("// thin: lifecycle DRIVEN by FireFly",CM),
 ("function mintOnSettlement(",KW),
 ("    bytes32 subId, address to, uint256 u)",TXTC),
 ("  onlyRole(ISSUER_ROLE) {",TXTC),
 ("  // RuleEngine must pass: eligible+juris",CM),
 ("  _mint(to, u);  emit Minted(subId,to,u);",TXTC),
 ("}",TXTC),
 ("function observeBarrier(uint256 lvlBps)",KW),
 ("  onlyRole(ORACLE_ROLE) { ",TXTC),
 ("  if (lvlBps < barrierBps) breached=true; }",TXTC),
 ("function settle(address h, uint256 u)",KW),
 ("  onlyRole(SETTLEMENT_ROLE) {",TXTC),
 ("  bool physical = breached; _burn(h,u); }",TXTC),
],ACCENT)
code(6.83,5.6,"orchestrator.ts  (FireFly SDK)",[
 ("// idempotent on subscriptionId",CM),
 ("async function onPaymentFinality(e){",KW),
 ("  if (await alreadyMinted(e.subId))",TXTC),
 ("      return;              // dedupe",CM),
 ("  const ok = await compliance",TXTC),
 ("      .check(e.wallet, e.subId);",TXTC),
 ("  if (!ok) return refundEscrow(e.subId);",TXTC),
 ("  await firefly.mintTokens(",TXTC),
 ("    { pool: POOL, to: e.wallet,",TXTC),
 ("      amount: e.units },",TXTC),
 ("    { requestId: e.subId,      // idem",CM),
 ("      confirm: true });   // wait final",CM),
 ("}",TXTC),
],TEAL)

# ---------- 10 RELIABILITY / THREATS ----------
s = slide(); header(s, "Reliability & threat model", "Failure-first design"); pagenum(s,10)
box(s,0.9,2.0,5.6,4.7,fill=WHITE,line=LINE)
txt(s,1.15,2.18,5.1,0.4,[[("Failure handling",16,ACCENT,True)]])
bullet(s,1.15,2.8,5.1,[
 ("Payment event lost", " → durable FireFly stream + re-drive"),
 ("Duplicate event", " → idempotency key, exactly-once mint"),
 ("Cash settled, mint reverts", " → auto-refund escrow, no partial state"),
 ("Chain reorg", " → confirm:true, only finalized receipts"),
 ("Oracle dispute at barrier", " → signed feed + Pause if systemic"),
],sz=13)
box(s,6.7,2.0,5.73,4.7,fill=WHITE,line=LINE)
txt(s,6.95,2.18,5.2,0.4,[[("Threats → controls",16,RGBColor(0xC0,0x3A,0x3A),True)]])
bullet(s,6.95,2.8,5.2,[
 ("ISSUER_ROLE key compromise", " → HSM/MPC, multisig, per-op limits"),
 ("Oracle manipulation", " → multi-source signed feed, dispute window"),
 ("Payment replay", " → end-to-end idempotency key"),
 ("Eligibility bypass on resale", " → RuleEngine gates every transfer"),
 ("PII on-chain", " → hashes only, private data exchange"),
 ("Reg S → US-person leakage", " → jurisdiction attr enforced on transfer"),
],sz=13,mk=RGBColor(0xC0,0x3A,0x3A))
txt(s,0.9,6.85,11.5,0.4,[[("Invariant: ",13,GOLD,True),("the reconciliation store — not the chain, not core banking alone — is the correlated source of truth. Both legs reconcile to it.",13,INK,False)]])

# ---------- 11 BUILD SEQUENCE ----------
s = slide(); header(s, "Pragmatic build sequence", "From zero to issuance"); pagenum(s,11)
steps=[
 ("1","Stand up FireFly","blockchain + token connector + event streams; prove mint/transfer/burn"),
 ("2","Deploy CMTAT base","permissioned EVM; wire RuleEngine with a stub eligibility rule"),
 ("3","Recon + ISO 20022","idempotency store + pacs.002 finality on ONE rail (SEPA Inst sandbox)"),
 ("4","Escrow mint-on-finality","prove exactly-once + auto-refund reversal paths"),
 ("5","Lifecycle","snapshot→coupon, barrier oracle, maturity settle (cash & physical)"),
 ("6","Compliance layer","KYC/AML, Travel Rule payloads, jurisdiction gating; PII off-chain"),
 ("7","Second jurisdiction","generalize matrix into RuleEngine config"),
 ("8","Audit + legal opinion","Solidity + orchestration audit, per-jurisdiction opinion before live"),
]
colw=5.75
for k,(num,t,d) in enumerate(steps):
    col=k//4; row=k%4
    x=0.9+col*6.0; y=2.05+row*1.15
    box(s,x,y,colw,1.0,fill=WHITE,line=LINE)
    c=box(s,x+0.14,y+0.19,0.62,0.62,fill=ACCENT if col==0 else TEAL,shape=MSO_SHAPE.OVAL)
    txt(s,x+0.14,y+0.19,0.62,0.62,[[(num,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.95,y+0.14,colw-1.05,0.4,[[(t,14.5,INK,True)]])
    txt(s,x+0.95,y+0.5,colw-1.05,0.45,[[(d,11,MUTE,False)]],sp=0.98)
txt(s,0.9,6.9,11.5,0.4,[[("No warranty. Validate every control with licensed counsel and your regulators before production issuance.",11.5,MUTE,False)]])

# ---------- 12 CLOSE ----------
s = slide(INK)
box(s,0,0,SW,SH,fill=INK,shape=MSO_SHAPE.RECTANGLE)
box(s,0.9,2.0,0.09,2.0,fill=GOLD,shape=MSO_SHAPE.RECTANGLE)
txt(s,1.2,2.15,11.2,1.9,[
 [("Fiat in. Fiat out. TradFi both ends.",32,WHITE,True)],
 [("The DLT is the shared message bus —",29,RGBColor(0xBF,0xD2,0xF0),True)],
 [("correlation, not value transfer, is the architecture.",29,RGBColor(0xBF,0xD2,0xF0),True)],
],sp=1.06)
txt(s,1.2,4.6,11,1.2,[
 [("• Recon store = source of truth   • Idempotency = exactly-once mint",16,WHITE,False)],
 [("• RuleEngine gates every transfer   • PII off-chain, hashes on   • Escrow + auto-reversal",16,WHITE,False)],
],sp=1.35)
txt(s,1.2,6.5,11,0.5,[[("Full write-up, diagrams & code skeletons in the companion repo & gist.  Not legal advice.",13,RGBColor(0x9A,0xAD,0xC8),False)]])

prs.save("RCN-CMTAT-FireFly-Gateway.pptx")
print("saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
